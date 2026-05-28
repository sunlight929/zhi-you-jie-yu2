"""
生成答辩展示 PPT（7 页 · 天大蓝白校徽风）
=====================================
精简版结构：
  1  封面（校徽 + 项目名）
  2  项目背景（1 页综合：数据 + 痛点 + 项目定位）
  3  设计思路（1 页综合：多源数据 + 技术路线）
  4  ▶  进入平台演示（重点：现场讲解作品）
  5  创新点（5 大亮点）
  6  团队分工 + 双导师
  7  致谢（校徽 + 团队署名）
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- 配色：天大蓝白 ----------
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_BG     = RGBColor(0xF4, 0xF8, 0xFC)
TJU_BLUE    = RGBColor(0x00, 0x3F, 0x7E)   # 天大蓝（主色）
TJU_DEEP    = RGBColor(0x00, 0x2A, 0x5C)   # 深蓝
TJU_LIGHT   = RGBColor(0x4A, 0x7B, 0xC0)   # 浅蓝
TJU_PALE    = RGBColor(0xDB, 0xE6, 0xF5)   # 极浅蓝底纹
TJU_GOLD    = RGBColor(0xC4, 0x96, 0x1A)   # 校徽辅助金（不抢戏）
TEXT_DARK   = RGBColor(0x1A, 0x28, 0x48)
TEXT_BODY   = RGBColor(0x33, 0x44, 0x66)
TEXT_GRAY   = RGBColor(0x70, 0x7A, 0x8E)
LIGHT_GRAY  = RGBColor(0xC8, 0xD2, 0xE0)

FONT = "PingFang SC"

ROOT = Path(__file__).resolve().parents[1]
SHOT_DIR = ROOT / "assets" / "screenshots"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height


# ============================================================
# 工具函数
# ============================================================
def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    return bg


def add_text(slide, text, left, top, width, height,
             font_size=14, color=TEXT_BODY, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font_name=FONT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_box(slide, left, top, width, height,
            fill=WHITE, border=TJU_BLUE, border_w=0.75, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, left, top, width, height)
    if rounded:
        box.adjustments[0] = 0.06
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if border:
        box.line.color.rgb = border
        box.line.width = Pt(border_w)
    else:
        box.line.fill.background()
    return box


def add_line_h(slide, left, top, width, color=TJU_BLUE, weight=2):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(weight))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_emblem_placeholder(slide, left, top, size=Inches(1.4)):
    """校徽占位：圆形蓝色边框 + 内部文字提示。
    用户打开 PPT 后，把真实校徽图片拖到这个位置即可替换。
    """
    # 外圈
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = WHITE
    circle.line.color.rgb = TJU_BLUE
    circle.line.width = Pt(2.5)

    # 内圈装饰
    inner_size = size - Inches(0.2)
    offset = Inches(0.1)
    inner = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    left + offset, top + offset,
                                    inner_size, inner_size)
    inner.fill.solid()
    inner.fill.fore_color.rgb = TJU_PALE
    inner.line.color.rgb = TJU_BLUE
    inner.line.width = Pt(0.5)

    # 中央文字
    add_text(slide, "校徽", left, top, size, size,
             font_size=12, color=TJU_BLUE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, "TJU", left, top + Inches(0.05), size, size,
             font_size=8, color=TJU_BLUE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)


def add_page_header(slide, section_label: str, page_title: str):
    """页面顶部：左侧蓝色色块 + 章节 + 标题。"""
    add_box(slide, Inches(0.5), Inches(0.45),
            Pt(5), Inches(0.6),
            fill=TJU_BLUE, border=None)
    add_text(slide, section_label, Inches(0.65), Inches(0.4),
             Inches(2.5), Inches(0.4),
             font_size=12, color=TJU_BLUE, bold=True)
    add_text(slide, page_title, Inches(0.65), Inches(0.7),
             Inches(11.5), Inches(0.5),
             font_size=22, color=TEXT_DARK, bold=True)
    add_line_h(slide, Inches(0.5), Inches(1.3),
              Inches(12.33), color=TJU_BLUE, weight=1.5)


def add_footer(slide, page_num: int, total: int = 7,
               text_color=TEXT_GRAY):
    add_text(slide, "知忧·解郁  ·  抑郁症风险识别与决策支持平台  ·  天津大学医学院",
             Inches(0.5), H - Inches(0.35),
             Inches(9), Inches(0.25),
             font_size=9, color=text_color)
    add_text(slide, f"{page_num} / {total}",
             W - Inches(1.5), H - Inches(0.35),
             Inches(1), Inches(0.25),
             font_size=9, color=text_color, align=PP_ALIGN.RIGHT)


# ============================================================
# Slide 1: 封面（校徽 + 项目名）
# ============================================================
def slide_1():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)

    # 顶部蓝色装饰条
    add_box(s, 0, 0, W, Inches(0.4), fill=TJU_BLUE, border=None)
    # 底部蓝色装饰条
    add_box(s, 0, H - Inches(0.4), W, Inches(0.4),
            fill=TJU_BLUE, border=None)

    # 校徽（左上角）
    add_emblem_placeholder(s, Inches(0.7), Inches(0.85), size=Inches(1.2))
    # 校徽（右上角）
    add_emblem_placeholder(s, W - Inches(1.9), Inches(0.85), size=Inches(1.2))

    # 校名
    add_text(s, "天津大学",
             Inches(2.0), Inches(0.95), Inches(2.5), Inches(0.5),
             font_size=18, color=TJU_BLUE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "TIANJIN UNIVERSITY",
             Inches(2.0), Inches(1.4), Inches(3.5), Inches(0.4),
             font_size=10, color=TJU_LIGHT,
             anchor=MSO_ANCHOR.TOP)
    add_text(s, "医学院",
             Inches(2.0), Inches(1.6), Inches(2.5), Inches(0.4),
             font_size=12, color=TEXT_GRAY,
             anchor=MSO_ANCHOR.TOP)

    # 大赛 LOGO 文字（右上）
    add_text(s, "中国大学生计算机设计大赛",
             W - Inches(5.5), Inches(1.0), Inches(3.5), Inches(0.5),
             font_size=12, color=TJU_BLUE, bold=True,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "Chinese Collegiate Computing Competition  ·  4C",
             W - Inches(5.5), Inches(1.4), Inches(3.5), Inches(0.4),
             font_size=9, color=TJU_LIGHT,
             align=PP_ALIGN.RIGHT)

    # 中央装饰金色横线
    add_line_h(s, Inches(5.665), Inches(2.6), Inches(2),
              color=TJU_GOLD, weight=3)

    # 主标题
    add_text(s, "知忧·解郁",
             Inches(1), Inches(2.95), Inches(11.33), Inches(1.4),
             font_size=64, color=TJU_DEEP, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, "多源数据驱动的抑郁症风险识别与决策支持平台",
             Inches(1), Inches(4.35), Inches(11.33), Inches(0.6),
             font_size=20, color=TJU_BLUE, bold=True,
             align=PP_ALIGN.CENTER)

    add_line_h(s, Inches(5.165), Inches(5.15), Inches(3),
              color=LIGHT_GRAY, weight=1)

    # 大赛信息
    add_text(s, "2026 中国大学生计算机设计大赛  ·  大数据应用赛道  ·  生物与医疗大数据小类",
             Inches(1), Inches(5.4), Inches(11.33), Inches(0.4),
             font_size=14, color=TEXT_BODY, align=PP_ALIGN.CENTER)
    add_text(s, "国赛参赛作品  ·  天津大学医学院",
             Inches(1), Inches(5.85), Inches(11.33), Inches(0.4),
             font_size=14, color=TEXT_BODY, align=PP_ALIGN.CENTER)

    # 团队
    add_text(s, "指导教师：张小臣 副研究员（医学方向）   /   张淑芳 副教授（AI 方向）",
             Inches(1), Inches(6.55), Inches(11.33), Inches(0.4),
             font_size=12, color=TJU_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, "答辩展示  ·  2026.05",
             Inches(1), Inches(6.95), Inches(11.33), Inches(0.3),
             font_size=10, color=TEXT_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 2: 项目背景（综合 1 页：数据 + 项目定位）
# ============================================================
def slide_2():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    add_page_header(s, "01  项目背景",
                   "我国精神卫生体系挑战  +  本项目定位")

    # 左半：核心数据
    add_text(s, "▎  我国精神卫生现状",
             Inches(0.7), Inches(1.55), Inches(6), Inches(0.4),
             font_size=14, color=TJU_BLUE, bold=True)

    stats = [
        ("9,500 万", "我国抑郁症患者总数",
         "Lu J 2021 Lancet Psychiatry"),
        ("24.7%", "大学生抑郁检出率",
         "Gao L 2020 Sci Rep Meta 分析"),
        ("3.4 / 10万", "精神科医师密度",
         "远低于发达国家 10-15"),
    ]
    for i, (val, label, src) in enumerate(stats):
        y = Inches(2.05 + i * 1.3)
        add_box(s, Inches(0.7), y, Inches(6), Inches(1.15),
                fill=SOFT_BG, border=TJU_BLUE, border_w=0.75, rounded=True)
        add_line_h(s, Inches(0.7) + Pt(2), y + Pt(2),
                   Inches(6) - Pt(4), color=TJU_BLUE, weight=2)
        # 大数字
        add_text(s, val, Inches(0.85), y + Inches(0.15),
                 Inches(2.5), Inches(0.85),
                 font_size=28, color=TJU_DEEP, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
        # 标签 + 来源
        add_text(s, label, Inches(3.4), y + Inches(0.15),
                 Inches(3.2), Inches(0.5),
                 font_size=13, color=TEXT_DARK, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, src, Inches(3.4), y + Inches(0.6),
                 Inches(3.2), Inches(0.4),
                 font_size=10, color=TEXT_GRAY, italic=True,
                 anchor=MSO_ANCHOR.MIDDLE)

    # 中部分隔
    sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(6.835), Inches(1.55),
                              Pt(1), Inches(5.0))
    sep.fill.solid()
    sep.fill.fore_color.rgb = LIGHT_GRAY
    sep.line.fill.background()

    # 右半：项目定位
    add_text(s, "▎  本项目定位",
             Inches(7.1), Inches(1.55), Inches(6), Inches(0.4),
             font_size=14, color=TJU_BLUE, bold=True)

    add_box(s, Inches(7.1), Inches(2.05), Inches(5.83),
            Inches(2.0), fill=TJU_PALE, border=TJU_BLUE,
            border_w=1.5, rounded=True)
    add_text(s, "面向我国大学生与中老年两大重点人群",
             Inches(7.3), Inches(2.2), Inches(5.5), Inches(0.4),
             font_size=12, color=TEXT_DARK, bold=True)
    add_text(s, "「真实学术公开数据 + 中国本土流行病学仿真」",
             Inches(7.3), Inches(2.6), Inches(5.5), Inches(0.4),
             font_size=12, color=TJU_BLUE, bold=True)
    add_text(s, "多源数据驱动的抑郁风险识别与决策支持平台",
             Inches(7.3), Inches(3.0), Inches(5.5), Inches(0.4),
             font_size=12, color=TEXT_DARK, bold=True)
    add_text(s, "结合 PHQ-9 / CES-D 国际临床金标准量表",
             Inches(7.3), Inches(3.45), Inches(5.5), Inches(0.4),
             font_size=11, color=TEXT_BODY)
    add_text(s, "+ 三模型机器学习集成 + 可解释性输出",
             Inches(7.3), Inches(3.75), Inches(5.5), Inches(0.4),
             font_size=11, color=TEXT_BODY)

    # 三大挑战
    add_text(s, "▎  三大结构性挑战",
             Inches(7.1), Inches(4.3), Inches(6), Inches(0.4),
             font_size=14, color=TJU_BLUE, bold=True)
    challenges = [
        ("识别难", "临床就诊率 < 10%"),
        ("资源缺", "心理治疗师严重短缺"),
        ("分布不均", "农村基层服务可及性差"),
    ]
    for i, (title, desc) in enumerate(challenges):
        y = Inches(4.8 + i * 0.55)
        add_box(s, Inches(7.1), y, Inches(5.83), Inches(0.45),
                fill=WHITE, border=TJU_GOLD, border_w=0.75)
        add_text(s, title, Inches(7.3), y + Inches(0.05),
                 Inches(1.5), Inches(0.35),
                 font_size=12, color=TJU_GOLD, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, desc, Inches(8.9), y + Inches(0.05),
                 Inches(4), Inches(0.35),
                 font_size=11, color=TEXT_BODY,
                 anchor=MSO_ANCHOR.MIDDLE)

    # 学术依托
    add_text(s, "学术依托：张小臣副研究员「面向脑认知的抑郁情绪识别模型开发」横向项目",
             Inches(7.1), Inches(6.7), Inches(5.83), Inches(0.3),
             font_size=10, color=TJU_BLUE, italic=True, bold=True)

    add_footer(s, 2)


# ============================================================
# Slide 3: 设计思路（综合 1 页：多源数据 + 技术路线）
# ============================================================
def slide_3():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    add_page_header(s, "02  设计思路",
                   "多源数据生态  +  技术路线")

    # 上方：多源数据
    add_text(s, "▎  多源数据生态  ·  共 14,750 条",
             Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
             font_size=14, color=TJU_BLUE, bold=True)

    sources = [
        ("🇺🇸  NHANES", "5,068", "美国 CDC", "Public Domain", TJU_BLUE),
        ("🌏  Mendeley", "682", "学生 PHQ-9", "CC BY 4.0", TJU_BLUE),
        ("🇨🇳  中国学生", "5,000", "Gao 2020 仿真", "Sci Rep", TJU_GOLD),
        ("🇨🇳  中老年", "4,000", "CHARLS 仿真", "北大", TJU_GOLD),
    ]
    card_w = Inches(2.95)
    gap = Inches(0.13)
    for i, (title, n, desc, lic, color) in enumerate(sources):
        x = Inches(0.7) + (card_w + gap) * i
        add_box(s, x, Inches(2.05), card_w, Inches(1.3),
                fill=SOFT_BG, border=color, border_w=1.0, rounded=True)
        add_line_h(s, x + Pt(2), Inches(2.05) + Pt(2),
                   card_w - Pt(4), color=color, weight=2)
        add_text(s, title, x, Inches(2.18),
                 card_w, Inches(0.4),
                 font_size=12, color=TEXT_DARK, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(s, n, x, Inches(2.55),
                 card_w, Inches(0.5),
                 font_size=24, color=color, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(s, desc, x, Inches(3.0),
                 card_w, Inches(0.3),
                 font_size=10, color=TEXT_BODY, align=PP_ALIGN.CENTER)
        add_text(s, lic, x, Inches(3.1),
                 card_w, Inches(0.3),
                 font_size=9, color=TEXT_GRAY, italic=True,
                 align=PP_ALIGN.CENTER)

    # 量表
    add_text(s, "📋  量表标准：PHQ-9（Kroenke 2001 · 学生 / 成人）  +  CES-D 10（Andresen 1994 · CHARLS 中老年标准）",
             Inches(0.7), Inches(3.55), Inches(12), Inches(0.4),
             font_size=11, color=TEXT_BODY)

    # 下方：技术路线
    add_text(s, "▎  技术路线  ·  四层架构 + 三模型集成 + 量表融合",
             Inches(0.7), Inches(4.1), Inches(12), Inches(0.4),
             font_size=14, color=TJU_BLUE, bold=True)

    # 四层架构紧凑版
    layers = [
        ("用户层", "个人 / 高校心理中心 / 社区卫生 / 公共卫生决策"),
        ("应用层", "Streamlit 5 大模块：驾驶舱 / 数据全景 / 风险分析 / 智能评估 / 天津决策"),
        ("模型层", "RF + GBDT + LR  →  量表融合（60% 量表 + 40% ML）"),
        ("数据层", "NHANES + Mendeley + 中国学生 + 中老年 = 14,750 条"),
    ]
    for i, (name, content) in enumerate(layers):
        y = Inches(4.6 + i * 0.5)
        add_box(s, Inches(0.7), y, Inches(1.6), Inches(0.42),
                fill=TJU_BLUE, border=None)
        add_text(s, name, Inches(0.7), y, Inches(1.6), Inches(0.42),
                 font_size=12, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_box(s, Inches(2.3), y, Inches(10.33), Inches(0.42),
                fill=SOFT_BG, border=TJU_BLUE, border_w=0.5)
        add_text(s, content, Inches(2.45), y, Inches(10.05), Inches(0.42),
                 font_size=11, color=TEXT_BODY,
                 anchor=MSO_ANCHOR.MIDDLE)

    # 关键性能
    add_text(s, "📊  内部验证集 AUC = 0.935  ·  召回率 0.866  ·  可解释性：Contribution = β × z",
             Inches(0.7), Inches(6.7), Inches(12), Inches(0.3),
             font_size=10, color=TJU_GOLD, italic=True, bold=True,
             align=PP_ALIGN.CENTER)

    add_footer(s, 3)


# ============================================================
# Slide 4: ▶ 进入平台演示（重点页）
# ============================================================
def slide_4():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)

    # 顶部蓝色装饰条
    add_box(s, 0, 0, W, Inches(0.4), fill=TJU_BLUE, border=None)
    add_box(s, 0, H - Inches(0.4), W, Inches(0.4),
            fill=TJU_BLUE, border=None)

    # 大标题
    add_text(s, "▶  进入平台实操演示",
             Inches(1), Inches(1.5), Inches(11.33), Inches(1.2),
             font_size=54, color=TJU_DEEP, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_line_h(s, Inches(5.165), Inches(2.85), Inches(3),
              color=TJU_GOLD, weight=3)

    # 链接（带超链接）
    link_box = add_box(s, Inches(3.165), Inches(3.3),
                       Inches(7), Inches(1.0),
                       fill=TJU_BLUE, border=None, rounded=True)
    tb = s.shapes.add_textbox(Inches(3.165), Inches(3.3),
                               Inches(7), Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "🔗   http://localhost:8501"
    run.font.name = "PingFang SC"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.hyperlink.address = "http://localhost:8501"

    # 5 个模块图标
    add_text(s, "本平台共五大功能模块",
             Inches(1), Inches(4.7), Inches(11.33), Inches(0.4),
             font_size=14, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    modules = [
        ("📡", "数据驾驶舱"),
        ("📊", "数据全景"),
        ("🔍", "风险因素分析"),
        ("🧠", "智能评估"),
        ("🏙️", "天津决策支持"),
    ]
    card_w = Inches(2.0)
    gap = Inches(0.2)
    total_w = card_w * 5 + gap * 4
    start_x = (W - total_w) / 2
    top = Inches(5.25)
    for i, (icon, name) in enumerate(modules):
        x = start_x + (card_w + gap) * i
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    x + (card_w - Inches(0.9)) / 2,
                                    top, Inches(0.9), Inches(0.9))
        circle.fill.solid()
        circle.fill.fore_color.rgb = TJU_PALE
        circle.line.color.rgb = TJU_BLUE
        circle.line.width = Pt(1)
        add_text(s, icon, x, top + Inches(0.12),
                 card_w, Inches(0.65),
                 font_size=28, align=PP_ALIGN.CENTER)
        add_text(s, name, x, top + Inches(1.05),
                 card_w, Inches(0.4),
                 font_size=12, color=TJU_DEEP, bold=True,
                 align=PP_ALIGN.CENTER)

    add_text(s, "下面切换到平台浏览器，进行 5 分钟实操讲解 →",
             Inches(1), Inches(6.85), Inches(11.33), Inches(0.4),
             font_size=14, color=TJU_GOLD, bold=True, italic=True,
             align=PP_ALIGN.CENTER)


# ============================================================
# Slide 5: 创新点
# ============================================================
def slide_5():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    add_page_header(s, "03  创新点",
                   "五大创新点  ·  方法学严谨 + 临床安全 + 本地化")

    innovations = [
        ("01", "多源数据生态",
         "真实公开（NHANES + Mendeley）+ 中国本土流行病学仿真双轨\n"
         "每个数据源透明标注真实/仿真徽章，不掩盖、不混淆"),
        ("02", "临床金标准量表完整集成",
         "PHQ-9（Kroenke 2001）+ CES-D 10（CHARLS 标准）原题原选项\n"
         "PHQ-9 第 9 题（自伤念头）独立预警机制，临床安全把关"),
        ("03", "三模型集成 + 量表临床切点融合",
         "RF + GBDT + LR 集成 AUC 0.935；量表 60% + ML 40% 加权\n"
         "保证 PHQ-9 = 0 必为低风险，临床主导，AI 辅助"),
        ("04", "局部贡献可解释性",
         "Contribution = β × z 轻量级分解\n"
         "每位用户的评估都附带「为什么风险高/低」的可视化解释"),
        ("05", "中国 34 省地图 + 天津 16 区县本地化",
         "ECharts 5.4.3 全国省级 Choropleth 热力图\n"
         "区县级风险分级 + 4 类干预建议 + 本地求助资源整合"),
    ]
    for i, (num, title, desc) in enumerate(innovations):
        y = Inches(1.6 + i * 1.05)
        add_box(s, Inches(0.7), y, Inches(0.85), Inches(0.95),
                fill=TJU_BLUE, border=None)
        add_text(s, num, Inches(0.7), y, Inches(0.85), Inches(0.95),
                 font_size=22, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_box(s, Inches(1.55), y, Inches(11.08), Inches(0.95),
                fill=SOFT_BG, border=TJU_BLUE, border_w=0.5)
        add_text(s, title, Inches(1.75), y + Inches(0.08),
                 Inches(10.7), Inches(0.4),
                 font_size=13, color=TJU_DEEP, bold=True)
        add_text(s, desc, Inches(1.75), y + Inches(0.42),
                 Inches(10.7), Inches(0.55),
                 font_size=10, color=TEXT_BODY)

    add_footer(s, 5)


# ============================================================
# Slide 6: 团队分工
# ============================================================
def slide_6():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)
    add_page_header(s, "04  团队分工",
                   "医工交叉双导师协同")

    # 学生
    add_text(s, "▎  学生作者（天津大学医学院  ·  临床医学专业本科生）",
             Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
             font_size=14, color=TJU_BLUE, bold=True)

    # 队长
    add_box(s, Inches(0.7), Inches(2.05), Inches(5.95), Inches(1.5),
            fill=SOFT_BG, border=TJU_BLUE, border_w=1.0, rounded=True)
    add_text(s, "👤  队长", Inches(0.9), Inches(2.2),
             Inches(5.6), Inches(0.4),
             font_size=15, color=TJU_BLUE, bold=True)
    add_text(s, "选题立项 / 医学专业表达 / 量表选型",
             Inches(0.9), Inches(2.65), Inches(5.6), Inches(0.4),
             font_size=12, color=TEXT_BODY)
    add_text(s, "研究报告撰写 / 模型训练 / Web 平台开发",
             Inches(0.9), Inches(3.05), Inches(5.6), Inches(0.4),
             font_size=12, color=TEXT_BODY)

    # 队员
    add_box(s, Inches(6.85), Inches(2.05), Inches(5.95), Inches(1.5),
            fill=SOFT_BG, border=TJU_LIGHT, border_w=1.0, rounded=True)
    add_text(s, "👤  队员", Inches(7.05), Inches(2.2),
             Inches(5.6), Inches(0.4),
             font_size=15, color=TJU_LIGHT, bold=True)
    add_text(s, "临床场景分析 / 风险因素解读",
             Inches(7.05), Inches(2.65), Inches(5.6), Inches(0.4),
             font_size=12, color=TEXT_BODY)
    add_text(s, "天津本地化策略 / 文档与演示",
             Inches(7.05), Inches(3.05), Inches(5.6), Inches(0.4),
             font_size=12, color=TEXT_BODY)

    # 双导师
    add_text(s, "▎  指导教师（双导师 · 医工交叉）",
             Inches(0.7), Inches(3.8), Inches(12), Inches(0.4),
             font_size=14, color=TJU_GOLD, bold=True)

    # 张小臣
    add_box(s, Inches(0.7), Inches(4.3), Inches(5.95), Inches(2.55),
            fill=SOFT_BG, border=TJU_GOLD, border_w=1.0, rounded=True)
    add_text(s, "🩺  张小臣  副研究员 / 正高级工程师",
             Inches(0.9), Inches(4.45), Inches(5.6), Inches(0.4),
             font_size=14, color=TJU_GOLD, bold=True)
    add_text(s, "天津大学医学院",
             Inches(0.9), Inches(4.85), Inches(5.6), Inches(0.4),
             font_size=11, color=TEXT_GRAY)
    add_text(s, "神经生物学  ·  抑郁症发病机制",
             Inches(0.9), Inches(5.25), Inches(5.6), Inches(0.4),
             font_size=11, color=TEXT_BODY)
    add_text(s, "「面向脑认知的抑郁情绪识别模型开发」项目负责人",
             Inches(0.9), Inches(5.65), Inches(5.6), Inches(0.4),
             font_size=11, color=TJU_BLUE, bold=True)
    add_text(s, "30+ SCI（Advanced Science / Neuron / PNAS）",
             Inches(0.9), Inches(6.05), Inches(5.6), Inches(0.4),
             font_size=11, color=TEXT_BODY)
    add_text(s, "→  医学专业方向把关",
             Inches(0.9), Inches(6.45), Inches(5.6), Inches(0.4),
             font_size=11, color=TJU_BLUE, italic=True, bold=True)

    # 张淑芳
    add_box(s, Inches(6.85), Inches(4.3), Inches(5.95), Inches(2.55),
            fill=SOFT_BG, border=TJU_GOLD, border_w=1.0, rounded=True)
    add_text(s, "🤖  张淑芳  副教授 / 博士生导师",
             Inches(7.05), Inches(4.45), Inches(5.6), Inches(0.4),
             font_size=14, color=TJU_GOLD, bold=True)
    add_text(s, "天津大学电气自动化与信息工程学院",
             Inches(7.05), Inches(4.85), Inches(5.6), Inches(0.4),
             font_size=11, color=TEXT_GRAY)
    add_text(s, "电子信息工程系  系主任",
             Inches(7.05), Inches(5.25), Inches(5.6), Inches(0.4),
             font_size=11, color=TEXT_GRAY)
    add_text(s, "人工智能  ·  AIGC 大模型  ·  数字图像处理",
             Inches(7.05), Inches(5.65), Inches(5.6), Inches(0.4),
             font_size=11, color=TEXT_BODY)
    add_text(s, "60+ SCI 论文  +  18 项授权国家发明专利",
             Inches(7.05), Inches(6.05), Inches(5.6), Inches(0.4),
             font_size=11, color=TEXT_BODY)
    add_text(s, "→  AI 技术方向把关",
             Inches(7.05), Inches(6.45), Inches(5.6), Inches(0.4),
             font_size=11, color=TJU_BLUE, italic=True, bold=True)

    add_footer(s, 6)


# ============================================================
# Slide 7: 致谢（校徽 + 团队署名）
# ============================================================
def slide_7():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, WHITE)

    # 顶部蓝色装饰条
    add_box(s, 0, 0, W, Inches(0.4), fill=TJU_BLUE, border=None)
    add_box(s, 0, H - Inches(0.4), W, Inches(0.4),
            fill=TJU_BLUE, border=None)

    # 校徽（顶部居中）
    add_emblem_placeholder(s, Inches(6.165), Inches(0.85),
                           size=Inches(1.0))

    add_line_h(s, Inches(5.665), Inches(2.1), Inches(2),
              color=TJU_GOLD, weight=3)

    # 引言
    add_text(s, "「识别一个抑郁症病人，可能就救一条命。」",
             Inches(1), Inches(2.5), Inches(11.33), Inches(0.7),
             font_size=26, color=TJU_DEEP, bold=True, italic=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "—— 张小臣 副研究员",
             Inches(1), Inches(3.3), Inches(11.33), Inches(0.4),
             font_size=12, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    # 寄语
    add_text(s, "我们希望这个平台",
             Inches(1), Inches(4.0), Inches(11.33), Inches(0.5),
             font_size=18, color=TEXT_BODY, align=PP_ALIGN.CENTER)
    add_text(s, "能成为更多医生、老师、家人手中的一个小工具",
             Inches(1), Inches(4.45), Inches(11.33), Inches(0.5),
             font_size=18, color=TEXT_BODY, align=PP_ALIGN.CENTER)
    add_text(s, "让被忧愁笼罩的人，早一点被看见",
             Inches(1), Inches(4.95), Inches(11.33), Inches(0.5),
             font_size=22, color=TJU_BLUE, bold=True, align=PP_ALIGN.CENTER)

    # 致谢
    add_line_h(s, Inches(5.165), Inches(5.95), Inches(3),
              color=LIGHT_GRAY, weight=1)
    add_text(s, "致  谢",
             Inches(1), Inches(6.05), Inches(11.33), Inches(0.4),
             font_size=14, color=TJU_DEEP, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "中国大学生计算机设计大赛组委会  ·  南开大学",
             Inches(1), Inches(6.4), Inches(11.33), Inches(0.4),
             font_size=12, color=TEXT_BODY, align=PP_ALIGN.CENTER)
    add_text(s, "张小臣副研究员  ·  张淑芳副教授",
             Inches(1), Inches(6.7), Inches(11.33), Inches(0.4),
             font_size=12, color=TJU_BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, "天津大学医学院  知忧解郁团队  ·  2026.05",
             Inches(1), Inches(7.0), Inches(11.33), Inches(0.3),
             font_size=10, color=TEXT_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# 生成
# ============================================================
slide_1()
slide_2()
slide_3()
slide_4()
slide_5()
slide_6()
slide_7()

OUT = Path(__file__).resolve().parents[2] / "知忧解郁_答辩展示_v1.0.pptx"
prs.save(OUT)
print(f"✅ PPT 已生成：{OUT}")
print(f"   共 {len(prs.slides)} 张幻灯片，文件大小 {OUT.stat().st_size // 1024} KB")
print()
print("配色：天大蓝白校徽风（#003F7E 主蓝 + #C4961A 辅金）")
print()
print("章节结构（7 张）：")
print("  1  封面（双校徽 + 项目名 + 双导师）")
print("  2  项目背景（左右合并：核心数据 + 项目定位）")
print("  3  设计思路（多源数据 + 技术路线 + AUC）")
print("  4  ▶ 进入平台演示（含 localhost:8501 超链接）— 重点 ✨")
print("  5  创新点（5 大）")
print("  6  团队分工（学生 + 双导师）")
print("  7  致谢（顶部校徽 + 张小臣老师引言）")
print()
print("📌 校徽占位说明：")
print("  打开 PPT 后，把真实天大校徽 PNG 图片拖到圆形占位上即可替换。")
